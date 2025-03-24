from PyPDF2 import PdfReader
from pptx import Presentation
from models import PitchDeckData
import os

class Parser:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[1].lower()

    def parse(self):
        if self.file_extension == '.pdf':
            return self._parse_pdf()
        elif self.file_extension == '.pptx':
            return self._parse_pptx()
        else:
            raise ValueError("Unsupported file format. Only PDF and PPTX files are allowed.")

    def _parse_pdf(self):
        try:
            reader = PdfReader(self.file_path)
            slides = []

            for page in reader.pages:
                text = page.extract_text()
                metadata = reader.metadata or {}
                slides.append({
                    'title': f"Slide {len(slides) + 1}",
                    'content': text or "No content",
                    'metadata': {
                        'author': metadata.get('/author', 'Unknown'),
                        'created_date': metadata.get('/creation_date', 'Unknown')
                    }
                })

            return slides
        except Exception as e:
            raise ValueError(f"Failed to parse PDF: {str(e)}")

    def _parse_pptx(self):
        try:
            presentation = Presentation(self.file_path)
            slides = []

            for i, slide in enumerate(presentation.slides):
                title = slide.shapes.title.text if slide.shapes.title else f"Slide {i + 1}"
                content = []

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)

                slides.append({
                    'title': title,
                    'content': "\n".join(content) or "No content",
                    'metadata': {
                        'author': presentation.core_properties.author or 'Unknown',
                        'created_date': presentation.core_properties.created or 'Unknown'
                    }
                })

            return slides
        except Exception as e:
            raise ValueError(f"Failed to parse PPTX: {str(e)}")


    def parse_and_store(self):
        parsed_data = self.parse()
        stored_count = 0

        for item in parsed_data:
            data = PitchDeckData(
                slide_title=item.get('slide_title', ''),
                text_content=item.get('text_content', ''),
                metadata=item.get('metadata', {})
            ).to_dict()
            PitchDeckData.insert_data(data)
            stored_count += 1

        return f"{stored_count} slides have been successfully stored in the database."

