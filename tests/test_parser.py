import unittest
import os
from app.parser import Parser

class TestParser(unittest.TestCase):
    def setUp(self):
        self.test_files_dir = os.path.join(os.path.dirname(__file__), 'test_files')
        if not os.path.exists(self.test_files_dir):
            os.makedirs(self.test_files_dir)

    def tearDown(self):
        for filename in os.listdir(self.test_files_dir):
            file_path = os.path.join(self.test_files_dir, filename)
            if os.path.isfile(file_path):
                os.unlink(file_path)

    def test_parse_pdf(self):
        pdf_content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /Contents 4 0 R /Resources << /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >> >>
endobj
4 0 obj
<< /Length 59 >>
stream
BT
/F1 12 Tf
72 720 Td
(Hello, World!) Tj
ET
endstream
endobj
xref
0 5
0000000000 65535 f
0000000010 00000 n
0000000079 00000 n
0000000138 00000 n
0000000217 00000 n
trailer
<< /Size 5 /Root 1 0 R >>
startxref
297
%%EOF
"""
        pdf_path = os.path.join(self.test_files_dir, 'test.pdf')
        with open(pdf_path, 'wb') as f:
            f.write(pdf_content)

        parser = Parser(pdf_path)
        result = parser.parse()

        self.assertIsInstance(result, list)
        self.assertTrue(len(result) > 0)
        self.assertEqual(result[0]['title'], "Slide 1")
        self.assertIn("Hello, World!", result[0]['content'])
        self.assertEqual(result[0]['metadata']['author'], "Unknown")
        self.assertEqual(result[0]['metadata']['created_date'], "Unknown")

    def test_parse_invalid_pdf(self):
        pdf_content = b"Invalid PDF content"
        pdf_path = os.path.join(self.test_files_dir, 'invalid.pdf')
        with open(pdf_path, 'wb') as f:
            f.write(pdf_content)

        parser = Parser(pdf_path)
        with self.assertRaises(ValueError):
            parser.parse()

    def test_parse_empty_pdf(self):
        pdf_path = os.path.join(self.test_files_dir, 'empty.pdf')
        with open(pdf_path, 'wb') as f:
            pass

        parser = Parser(pdf_path)
        with self.assertRaises(ValueError):
            parser.parse()

    def test_parse_pptx(self):
        from pptx import Presentation
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Test PPTX"
        pptx_path = os.path.join(self.test_files_dir, 'test.pptx')
        presentation.save(pptx_path)

        parser = Parser(pptx_path)
        result = parser.parse()

        self.assertIsInstance(result, list)
        self.assertTrue(len(result) > 0)
        self.assertEqual(result[0]['title'], "Test PPTX")

    def test_unsupported_format(self):
        txt_path = os.path.join(self.test_files_dir, 'test.txt')
        with open(txt_path, 'w') as f:
            f.write("This is a text file.")

        parser = Parser(txt_path)
        with self.assertRaises(ValueError):
            parser.parse()

if __name__ == '__main__':
    unittest.main()